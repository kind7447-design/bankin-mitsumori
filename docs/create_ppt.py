from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ============================================================
# 共通設定
# ============================================================
W = Inches(13.33)  # ワイドスクリーン 16:9
H = Inches(7.5)

BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)  # 濃紺
BG_LIGHT  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE      = RGBColor(0x21, 0x96, 0xF3)
ACCENT    = RGBColor(0x00, 0xBF, 0xD8)
TEXT_W    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_D    = RGBColor(0x22, 0x22, 0x22)
GRAY      = RGBColor(0xF5, 0xF5, 0xF5)
GREEN     = RGBColor(0x43, 0xA0, 0x47)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=24, bold=False, color=TEXT_W,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================
# ① トライアル顧客向け操作マニュアル
# ============================================================
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # blank

slides_trial = [
    # (タイトル大, タイトル小, ボディ行リスト, bg_dark)
    {
        "big":   "板金見積システム 操作ガイド",
        "sub":   "〜図面をアップするだけで見積が完成〜",
        "body":  [],
        "dark":  True,
        "note":  "このビデオでは、板金見積システムの基本的な使い方を4分でご説明します。",
    },
    {
        "big":   "操作はたった4ステップ",
        "sub":   "",
        "body":  [
            "Step 1　図面アップロード",
            "Step 2　仕様確認・修正",
            "Step 3　加工工程の選択",
            "Step 4　金額確認・Excel出力",
        ],
        "dark":  True,
        "note":  "図面をアップロードするだけで、AIが自動的に材質・サイズを読み取り、見積金額を計算します。",
    },
    {
        "big":   "Step 1　図面アップロード",
        "sub":   "PDF・画像・DXFに対応",
        "body":  [
            "► 図面をドラッグ＆ドロップするだけ",
            "► AIが図番・部品名・材質・板厚・製品サイズを自動読み取り",
            "► 手動入力も可能",
        ],
        "dark":  False,
        "note":  "まず図面をドラッグ＆ドロップします。AIが図面を解析し、各項目を自動で入力します。",
    },
    {
        "big":   "Step 2　仕様確認・修正",
        "sub":   "黄色ハイライト＝AI推定箇所",
        "body":  [
            "► 材質・板厚・素材サイズ・製品サイズを確認",
            "► 間違いがあればその場で修正",
            "► 数量もここで設定",
        ],
        "dark":  False,
        "note":  "AIが読み取った値を確認します。黄色くハイライトされている項目がAIが推定した箇所です。",
    },
    {
        "big":   "Step 3　加工工程",
        "sub":   "必要な工程にチェックを入れる",
        "body":  [
            "► レーザー切断・バリ取りは通常ON",
            "► 曲げ加工・タップ加工は数量を入力",
            "► 単価はマスタで自由に変更可能",
        ],
        "dark":  False,
        "note":  "加工工程を選択します。各工程の単価はマスタで自由に変更できます。",
    },
    {
        "big":   "Step 4　金額確認・出力",
        "sub":   "粗利調整 → 作成者選択 → Excel出力",
        "body":  [
            "► 粗利率スライダーでリアルタイム金額調整",
            "► Excel出力と同時に自動保存",
            "► PDF出力も可能",
        ],
        "dark":  False,
        "note":  "粗利率をスライダーで調整して、Excel出力ボタンを押すと見積書がダウンロードされます。",
    },
    {
        "big":   "複数部品をまとめて見積",
        "sub":   "「＋行を追加」で1枚にまとめられる",
        "body":  [
            "► 部品ごとに入力 → 「行を追加」",
            "► 全部品を入力したらまとめてExcel出力",
            "► 合計金額も自動計算",
        ],
        "dark":  False,
        "note":  "複数の部品をまとめて見積もることもできます。",
    },
    {
        "big":   "見積履歴",
        "sub":   "いつでも呼び出し・再利用",
        "body":  [
            "► 「履歴」タブから過去の見積を表示",
            "► 編集・複製・削除が可能",
            "► 同じ部品の再見積もりに便利",
        ],
        "dark":  False,
        "note":  "上部の履歴タブから過去の見積を呼び出せます。",
    },
    {
        "big":   "マスタ設定",
        "sub":   "材料単価・加工単価を管理",
        "body":  [
            "► 「マスタ」タブから変更",
            "► 変更は即座に全見積に反映",
            "► 作成者マスタも管理可能",
        ],
        "dark":  False,
        "note":  "材料単価や加工単価が変わった場合は、マスタタブから変更できます。",
    },
    {
        "big":   "まとめ",
        "sub":   "",
        "body":  [
            "✅  図面をアップするだけで自動見積",
            "✅  複数部品をまとめて1枚の見積書に",
            "✅  Excel・PDFで即出力",
            "✅  見積履歴をいつでも呼び出し",
        ],
        "dark":  True,
        "note":  "以上が基本的な操作の流れです。トライアル期間中、存分にお試しください。",
    },
]

for s in slides_trial:
    slide = prs.slides.add_slide(blank)
    dark = s["dark"]
    set_bg(slide, BG_DARK if dark else BG_LIGHT)
    txt_color = TEXT_W if dark else TEXT_D

    # アクセントバー（左端）
    add_rect(slide, Inches(0), Inches(0), Inches(0.12), H, BLUE)

    # タイトル大
    add_text_box(slide, s["big"],
                 Inches(0.3), Inches(0.4), Inches(12.5), Inches(1.2),
                 font_size=36, bold=True, color=TEXT_W if dark else BLUE,
                 align=PP_ALIGN.LEFT)

    # サブタイトル
    if s["sub"]:
        add_text_box(slide, s["sub"],
                     Inches(0.3), Inches(1.5), Inches(12), Inches(0.6),
                     font_size=20, bold=False,
                     color=ACCENT if dark else RGBColor(0x55,0x55,0x55),
                     align=PP_ALIGN.LEFT)

    # ボディ
    if s["body"]:
        y_start = Inches(2.2) if s["sub"] else Inches(1.8)
        line_h = Inches(0.7)
        for i, line in enumerate(s["body"]):
            add_text_box(slide, line,
                         Inches(0.5), y_start + line_h * i,
                         Inches(12), Inches(0.65),
                         font_size=22, bold=False,
                         color=txt_color,
                         align=PP_ALIGN.LEFT)

    # ノート
    slide.notes_slide.notes_text_frame.text = s["note"]

out1 = r"C:\Users\isobe\マイドライブ\claude\bankin-mitsumori\docs\操作マニュアル_トライアル顧客向け.pptx"
prs.save(out1)
print(f"saved: {out1}")


# ============================================================
# ② 営業デモ用
# ============================================================
prs2 = Presentation()
prs2.slide_width  = W
prs2.slide_height = H

BG_SALES  = RGBColor(0x0D, 0x47, 0xA1)   # 濃いブルー
YELLOW    = RGBColor(0xFF, 0xD6, 0x00)
LIGHT_BG  = RGBColor(0xF0, 0xF4, 0xFF)

slides_sales = [
    {
        "big":   "板金見積、\n今まで何分かかっていましたか？",
        "sub":   "",
        "body":  [],
        "style": "impact",
        "note":  "板金の見積もり、1件あたり何分かかっていますか？熟練者でも15〜30分かかります。",
    },
    {
        "big":   "こんなお悩みはありませんか？",
        "sub":   "",
        "body":  [
            "😓　見積に時間がかかりすぎる",
            "😓　担当者によって金額がバラつく",
            "😓　急ぎの見積に対応できない",
            "😓　ベテランが辞めたら見積が止まる",
        ],
        "style": "problem",
        "note":  "板金業の見積は属人化しがちです。こうした課題を抱える会社様は多いのではないでしょうか。",
    },
    {
        "big":   "図面をアップするだけで\n見積が完成します",
        "sub":   "AIが材料費・加工費を自動計算",
        "body":  [],
        "style": "solution",
        "note":  "板金見積システムは、図面をアップロードするだけで自動的に見積を作成します。",
    },
    {
        "big":   "▶ デモ",
        "sub":   "実際の操作をご覧ください",
        "body":  [
            "① 図面（PDF）をドロップ",
            "② AIが材質・サイズを自動読み取り",
            "③ 加工工程を確認",
            "④ Excelで見積書が完成　→　約1分",
        ],
        "style": "demo",
        "note":  "実際に見てみましょう。図面をドラッグ＆ドロップすると数秒でAIが解析します。この操作、約1分です。",
    },
    {
        "big":   "導入前 → 導入後",
        "sub":   "",
        "body":  [
            "見積時間：20分　→　1分",
            "必要スキル：熟練者のみ　→　誰でも",
            "金額のブレ：あり　→　なし（単価マスタで統一）",
            "履歴管理：紙・Excel　→　クラウドで自動保存",
        ],
        "style": "benefit",
        "note":  "見積時間が20分から1分に短縮。担当者に関わらず均一な品質の見積が作れます。",
    },
    {
        "big":   "まずは14日間、無料でお試しください",
        "sub":   "",
        "body":  [
            "✅　初期費用：なし",
            "✅　インストール不要（ブラウザで使用）",
            "✅　専用URLを発行してすぐ使える",
            "✅　設定はすべてこちらで対応",
        ],
        "style": "trial",
        "note":  "まずは14日間、無料でお試しいただけます。設定はこちらで行いますので、今日から始められます。",
    },
    {
        "big":   "お問い合わせ・お申し込み",
        "sub":   "無料トライアルはこちらから",
        "body":  [
            "",
            "　　連絡先：（ここに記入）",
            "　　メール：（ここに記入）",
        ],
        "style": "closing",
        "note":  "ご興味をお持ちの方は、ぜひ無料トライアルをお試しください。",
    },
]

for s in slides_sales:
    slide = prs2.slides.add_slide(blank)
    style = s["style"]

    # 背景
    if style in ("impact", "solution", "closing"):
        set_bg(slide, BG_SALES)
        title_color = YELLOW
        body_color  = TEXT_W
    elif style == "problem":
        set_bg(slide, RGBColor(0x18, 0x18, 0x18))
        title_color = RGBColor(0xFF, 0x57, 0x22)
        body_color  = TEXT_W
    elif style == "demo":
        set_bg(slide, BG_DARK)
        title_color = ACCENT
        body_color  = TEXT_W
    elif style == "benefit":
        set_bg(slide, LIGHT_BG)
        title_color = RGBColor(0x0D, 0x47, 0xA1)
        body_color  = TEXT_D
    elif style == "trial":
        set_bg(slide, RGBColor(0xE8, 0xF5, 0xE9))
        title_color = GREEN
        body_color  = TEXT_D
    else:
        set_bg(slide, BG_LIGHT)
        title_color = BLUE
        body_color  = TEXT_D

    # アクセントバー
    bar_color = YELLOW if style in ("impact","solution","closing") else BLUE
    add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, bar_color)

    # タイトル
    font_big = 40 if style == "impact" else 34
    add_text_box(slide, s["big"],
                 Inches(0.35), Inches(0.35), Inches(12.5), Inches(1.6),
                 font_size=font_big, bold=True, color=title_color,
                 align=PP_ALIGN.LEFT)

    # サブ
    if s["sub"]:
        add_text_box(slide, s["sub"],
                     Inches(0.35), Inches(1.9), Inches(12), Inches(0.6),
                     font_size=20, color=RGBColor(0xAA,0xAA,0xAA) if style in ("impact","demo","solution","closing") else RGBColor(0x55,0x55,0x55),
                     align=PP_ALIGN.LEFT)

    # ボディ
    if s["body"]:
        y0 = Inches(2.5) if s["sub"] else Inches(2.1)
        lh = Inches(0.78)
        for i, line in enumerate(s["body"]):
            add_text_box(slide, line,
                         Inches(0.5), y0 + lh * i,
                         Inches(12), Inches(0.72),
                         font_size=23, color=body_color,
                         align=PP_ALIGN.LEFT)

    slide.notes_slide.notes_text_frame.text = s["note"]

out2 = r"C:\Users\isobe\マイドライブ\claude\bankin-mitsumori\docs\営業デモ用スライド.pptx"
prs2.save(out2)
print(f"saved: {out2}")
